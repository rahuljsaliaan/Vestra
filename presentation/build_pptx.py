#!/usr/bin/env python3
"""Build OutfitBuddy.pptx — v2.

A native, editable PowerPoint deck that mirrors the HTML presentation:
 - radial-gradient dark backgrounds, gradient accent fills, glows, soft shadows
 - every slide plays an AUTOMATIC staggered entrance choreography
   (float-up, pop-with-overshoot, directional wipes, looping motion paths)
 - fade transitions between slides

Animations are authored as raw OOXML <p:timing> behaviour trees, which both
PowerPoint and LibreOffice Impress play back.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml

import os
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "OutfitBuddy.pptx")

# ---------------- palette ----------------
PAPER  = RGBColor(0xF3, 0xEB, 0xDD)
MUTED  = RGBColor(0xB3, 0xA6, 0x90)
FAINT  = RGBColor(0x8A, 0x7D, 0x69)
ACCENT = RGBColor(0xF2, 0x95, 0x3F)
DARKTX = RGBColor(0x1A, 0x12, 0x0A)
CARD   = RGBColor(0x26, 0x1F, 0x16)
CARD2  = RGBColor(0x2E, 0x26, 0x1B)
LINE   = RGBColor(0x45, 0x3B, 0x2D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
AMAZON = RGBColor(0xFF, 0x99, 0x00)
FLIPK  = RGBColor(0x4A, 0x9B, 0xFF)
MYNTRA = RGBColor(0xFF, 0x5F, 0x8F)
AJIO   = RGBColor(0x9C, 0xB4, 0xC4)
TATA   = RGBColor(0xFF, 0x6B, 0x63)

ACC_HEX, ACC2_HEX = "F2953F", "DF7226"

SERIF = "Georgia"
SANS  = "Calibri"

STORES = [("Amazon", AMAZON), ("Flipkart", FLIPK), ("Myntra", MYNTRA), ("Ajio", AJIO), ("Tata CLiQ", TATA)]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# =====================================================================
# low-level helpers
# =====================================================================

def _xml(s):
    return parse_xml(s)


BG_XML = (
    f'<p:bg {nsdecls("p", "a")}><p:bgPr><a:gradFill rotWithShape="1"><a:gsLst>'
    '<a:gs pos="0"><a:srgbClr val="2C2214"/></a:gs>'
    '<a:gs pos="50000"><a:srgbClr val="181410"/></a:gs>'
    '<a:gs pos="100000"><a:srgbClr val="110E09"/></a:gs></a:gsLst>'
    '<a:path path="circle"><a:fillToRect l="80000" t="0" r="20000" b="100000"/></a:path>'
    '</a:gradFill><a:effectLst/></p:bgPr></p:bg>'
)


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s._element.cSld.insert(0, _xml(BG_XML))
    return s


def _add_run(p, rd):
    r = p.add_run()
    r.text = rd["t"]
    f = r.font
    f.size = Pt(rd.get("size", 18))
    f.name = rd.get("font", SANS)
    f.bold = rd.get("bold", False)
    f.italic = rd.get("italic", False)
    f.color.rgb = rd.get("color", PAPER)
    if rd.get("spc"):
        r._r.get_or_add_rPr().set("spc", str(rd["spc"]))
    if rd.get("strike"):
        r._r.get_or_add_rPr().set("strike", "sngStrike")


def _apply_runs(tf, runs, align, anchor, line_spacing):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [{"t": runs}]
    for rd in runs:
        _add_run(p, rd)


def set_paras(tf, paras, anchor=MSO_ANCHOR.MIDDLE):
    """Multi-line text as REAL paragraphs (never \\n breaks — LibreOffice
    staggers centred lines that share a paragraph)."""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, pd in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pd.get("align", PP_ALIGN.CENTER)
        if pd.get("line_spacing"):
            p.line_spacing = pd["line_spacing"]
        if pd.get("space_before"):
            p.space_before = Pt(pd["space_before"])
        if pd.get("space_after"):
            p.space_after = Pt(pd["space_after"])
        for rd in pd["runs"]:
            _add_run(p, rd)


def text(shapes, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _apply_runs(tb.text_frame, runs, align, anchor, line_spacing)
    return tb


def rrect(shapes, l, t, w, h, fill=None, line=None, line_w=1.25, radius=0.14,
          shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def shape_text(sp, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=None):
    _apply_runs(sp.text_frame, runs, align, anchor, line_spacing)
    return sp


def oval(shapes, l, t, w, h, fill=ACCENT, line=None):
    sp = shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    sp.shadow.inherit = False
    return sp


def arrow(shapes, x1, y1, x2, y2, color=ACCENT, width=2.0, head=True):
    cn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    if head:
        ln = cn.line._get_or_add_ln()
        te = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(te)
    cn.shadow.inherit = False
    return cn


def set_gradient(sp, c1, c2, angle=135, alpha2=None):
    """Replace a shape's fill with a linear gradient c1→c2."""
    spPr = sp._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)
    a2 = f'<a:alpha val="{alpha2}"/>' if alpha2 is not None else ""
    grad = _xml(
        f'<a:gradFill {nsdecls("a")} rotWithShape="1"><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{c2}">{a2}</a:srgbClr></a:gs>'
        f'</a:gsLst><a:lin ang="{int(angle * 60000)}" scaled="1"/></a:gradFill>'
    )
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)


def _effect_lst(sp):
    spPr = sp._element.spPr
    el = spPr.find(qn("a:effectLst"))
    if el is None:
        el = _xml(f'<a:effectLst {nsdecls("a")}/>')
        spPr.append(el)
    return el


def add_glow(sp, hexcolor=ACC_HEX, radius_pt=14, alpha=42000):
    el = _effect_lst(sp)
    el.insert(0, _xml(  # glow must precede outerShdw in the schema sequence
        f'<a:glow {nsdecls("a")} rad="{int(radius_pt * 12700)}">'
        f'<a:srgbClr val="{hexcolor}"><a:alpha val="{alpha}"/></a:srgbClr></a:glow>'
    ))


def add_shadow(sp, blur_pt=14, dist_pt=5, alpha=40000):
    el = _effect_lst(sp)
    el.append(_xml(
        f'<a:outerShdw {nsdecls("a")} blurRad="{int(blur_pt * 12700)}" dist="{int(dist_pt * 12700)}" '
        f'dir="5400000" rotWithShape="0"><a:srgbClr val="000000"><a:alpha val="{alpha}"/></a:srgbClr></a:outerShdw>'
    ))


def ring(shapes, l, t, d, alpha=13000, width_pt=2.25):
    """Decorative accent circle outline, very low alpha."""
    sp = shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    sp.fill.background()
    sp.line.color.rgb = ACCENT
    sp.line.width = Pt(width_pt)
    srgb = sp.line._get_or_add_ln().find(".//" + qn("a:srgbClr"))
    srgb.append(_xml(f'<a:alpha {nsdecls("a")} val="{alpha}"/>'))
    sp.shadow.inherit = False
    return sp


def eyebrow(slide, s, l=0.92, t=0.6, w=11.5, align=PP_ALIGN.LEFT):
    return text(slide.shapes, l, t, w, 0.4,
                [{"t": s.upper(), "size": 13, "color": ACCENT, "bold": True, "spc": 280}], align=align)


def chrome(slide, num, total=9):
    """Small brand mark + slide counter, like the HTML deck's fixed chrome."""
    text(slide.shapes, 0.55, 0.22, 2.5, 0.35, [
        {"t": "Outfit Buddy", "size": 15, "color": FAINT, "font": SERIF, "bold": True},
        {"t": ".", "size": 15, "color": ACCENT, "font": SERIF, "bold": True},
    ])
    text(slide.shapes, 11.9, 7.02, 1.1, 0.33,
         [{"t": f"{num} / {total}", "size": 11, "color": FAINT}], align=PP_ALIGN.RIGHT)


# =====================================================================
# animation timeline (raw OOXML timing behaviours)
# =====================================================================

class Timeline:
    """Collects entrance/emphasis/path effects for one slide, all starting
    automatically (withEffect) at absolute millisecond delays, then writes the
    <p:transition> + <p:timing> parts."""

    def __init__(self):
        self.nodes = []
        self.blds = []
        self._grp = {}
        self._id = 4

    def nid(self):
        self._id += 1
        return self._id

    def _grpid(self, spid):
        g = self._grp.get(spid, -1) + 1
        self._grp[spid] = g
        self.blds.append((spid, g))
        return g

    def _wrap(self, spid, delay, klass, pid, psub, inner):
        g = self._grpid(spid)
        return (
            f'<p:par><p:cTn id="{self.nid()}" presetID="{pid}" presetClass="{klass}" '
            f'presetSubtype="{psub}" fill="hold" grpId="{g}" nodeType="withEffect">'
            f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
            f'<p:childTnLst>{inner}</p:childTnLst></p:cTn></p:par>'
        )

    def _set_visible(self, spid):
        return (
            f'<p:set><p:cBhvr><p:cTn id="{self.nid()}" dur="1" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
        )

    def _fade_b(self, spid, dur):
        return (
            f'<p:animEffect transition="in" filter="fade"><p:cBhvr>'
            f'<p:cTn id="{self.nid()}" dur="{dur}"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>'
        )

    # --- entrance effects -------------------------------------------------
    def fade(self, sp, delay, dur=500):
        spid = sp.shape_id
        self.nodes.append(self._wrap(spid, delay, "entr", 10, 0,
                                     self._set_visible(spid) + self._fade_b(spid, dur)))

    def rise(self, sp, delay, dur=650, dy=0.05):
        """Fade in while floating up from dy (fraction of slide height)."""
        spid = sp.shape_id
        y = (
            f'<p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base">'
            f'<p:cTn id="{self.nid()}" dur="{dur}" decel="65000" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst></p:cBhvr>'
            f'<p:tavLst><p:tav tm="0"><p:val><p:strVal val="#ppt_y+{dy}"/></p:val></p:tav>'
            f'<p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav></p:tavLst></p:anim>'
        )
        self.nodes.append(self._wrap(spid, delay, "entr", 2, 4,
                                     self._set_visible(spid) + self._fade_b(spid, int(dur * 0.8)) + y))

    def pop(self, sp, delay, dur=550):
        """Fade in while scaling 78% → 104% → 100% (spring overshoot)."""
        spid = sp.shape_id
        d1 = int(dur * 0.7)
        d2 = dur - d1
        s1 = (
            f'<p:animScale><p:cBhvr><p:cTn id="{self.nid()}" dur="{d1}" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
            f'<p:from x="78000" y="78000"/><p:to x="104000" y="104000"/></p:animScale>'
        )
        s2 = (
            f'<p:animScale><p:cBhvr><p:cTn id="{self.nid()}" dur="{d2}" fill="hold">'
            f'<p:stCondLst><p:cond delay="{d1}"/></p:stCondLst></p:cTn>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
            f'<p:from x="104000" y="104000"/><p:to x="100000" y="100000"/></p:animScale>'
        )
        self.nodes.append(self._wrap(spid, delay, "entr", 23, 0,
                                     self._set_visible(spid) + self._fade_b(spid, int(dur * 0.6)) + s1 + s2))

    def wipe(self, sp, delay, dur=450, direction="right"):
        """Directional reveal (wipe(right) = sweeps left→right)."""
        spid = sp.shape_id
        b = (
            f'<p:animEffect transition="in" filter="wipe({direction})"><p:cBhvr>'
            f'<p:cTn id="{self.nid()}" dur="{dur}"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>'
        )
        self.nodes.append(self._wrap(spid, delay, "entr", 22, 0, self._set_visible(spid) + b))

    def zoom_slow(self, sp, delay, dur=900):
        """Gentle zoom from 92% with fade — for big closing statements."""
        spid = sp.shape_id
        s1 = (
            f'<p:animScale><p:cBhvr><p:cTn id="{self.nid()}" dur="{dur}" decel="70000" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
            f'<p:from x="92000" y="92000"/><p:to x="100000" y="100000"/></p:animScale>'
        )
        self.nodes.append(self._wrap(spid, delay, "entr", 23, 0,
                                     self._set_visible(spid) + self._fade_b(spid, int(dur * 0.7)) + s1))

    # --- emphasis / motion -------------------------------------------------
    def pulse(self, sp, delay, dur=450, scale=106000):
        spid = sp.shape_id
        d1 = dur // 2
        s1 = (
            f'<p:animScale><p:cBhvr><p:cTn id="{self.nid()}" dur="{d1}" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
            f'<p:from x="100000" y="100000"/><p:to x="{scale}" y="{scale}"/></p:animScale>'
        )
        s2 = (
            f'<p:animScale><p:cBhvr><p:cTn id="{self.nid()}" dur="{dur - d1}" fill="hold">'
            f'<p:stCondLst><p:cond delay="{d1}"/></p:stCondLst></p:cTn>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
            f'<p:from x="{scale}" y="{scale}"/><p:to x="100000" y="100000"/></p:animScale>'
        )
        self.nodes.append(self._wrap(spid, delay, "emph", 26, 0, s1 + s2))

    def travel(self, sp, delay, dx, dur=1700):
        """Loop the shape rightwards by dx (fraction of slide width) forever."""
        spid = sp.shape_id
        b = (
            f'<p:animMotion origin="layout" path="M 0 0 L {dx:.4f} 0 E" pathEditMode="relative">'
            f'<p:cBhvr><p:cTn id="{self.nid()}" dur="{dur}" repeatCount="indefinite"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>ppt_x</p:attrName><p:attrName>ppt_y</p:attrName></p:attrNameLst>'
            f'</p:cBhvr></p:animMotion>'
        )
        self.nodes.append(self._wrap(spid, delay, "path", 0, 0, b))

    # --- write to slide ----------------------------------------------------
    def finish(self, slide):
        slide._element.append(_xml(f'<p:transition {nsdecls("p")} spd="med"><p:fade/></p:transition>'))
        inner = "".join(self.nodes)
        blds = "".join(f'<p:bldP spid="{s}" grpId="{g}"/>' for s, g in self.blds)
        slide._element.append(_xml(
            f'<p:timing {nsdecls("p", "a")}><p:tnLst><p:par>'
            f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
            f'<p:seq concurrent="1" nextAc="seek">'
            f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            f'<p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
            f'<p:par><p:cTn id="4" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
            f"{inner}"
            f'</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn>'
            f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>'
            f'<p:bldLst>{blds}</p:bldLst></p:timing>'
        ))
        return len(self.nodes)


counts = []

# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = new_slide()
tl = Timeline()
r1 = ring(s.shapes, 9.9, -1.6, 4.4)
r2 = ring(s.shapes, -1.3, 5.7, 3.0, alpha=9000)
eb = eyebrow(s, "Fashion, aggregated")
brand = text(s.shapes, 0.9, 1.28, 11.8, 2.15, [
    {"t": "Outfit Buddy", "size": 120, "color": PAPER, "font": SERIF, "bold": True},
    {"t": ".", "size": 120, "color": ACCENT, "font": SERIF, "bold": True},
])
tag = text(s.shapes, 0.92, 3.82, 10, 0.75, [{"t": "One closet. Every store.", "size": 34, "color": PAPER, "bold": True}])
bar = rrect(s.shapes, 0.95, 4.6, 3.7, 0.07, fill=ACCENT, radius=0, shape=MSO_SHAPE.RECTANGLE)
set_gradient(bar, ACC_HEX, ACC_HEX, angle=0, alpha2=0)
sub = text(s.shapes, 0.92, 4.88, 8.6, 1.1,
           [{"t": "Browse fashion in one place — and always buy at the best price across India's biggest stores.",
             "size": 19, "color": MUTED}], line_spacing=1.25)
chips = []
x = 0.92
for name, _c in STORES:
    w = 0.52 + len(name) * 0.145
    c = rrect(s.shapes, x, 6.05, w, 0.55, fill=None, line=LINE, radius=0.5)
    shape_text(c, [{"t": name, "size": 14, "color": MUTED, "bold": True}])
    chips.append(c)
    x += w + 0.22

tl.fade(r1, 300, 900); tl.fade(r2, 500, 900)
tl.fade(eb, 0, 400)
tl.rise(brand, 120, dur=750, dy=0.055)
tl.rise(tag, 430)
tl.wipe(bar, 700, dur=550)
tl.rise(sub, 860)
for i, c in enumerate(chips):
    tl.pop(c, 1120 + i * 110)
counts.append(tl.finish(s))

# =====================================================================
# SLIDE 2 — WHY / PROBLEM (browser tabs)
# =====================================================================
s = new_slide()
tl = Timeline()
chrome(s, 2)
eb = eyebrow(s, "The problem")
h1 = text(s.shapes, 0.9, 1.1, 11.6, 0.95, [{"t": "The same jacket,", "size": 46, "color": PAPER, "font": SERIF, "bold": True}])
h2 = text(s.shapes, 0.9, 1.98, 11.6, 0.95, [{"t": "five different prices.", "size": 46, "color": ACCENT, "font": SERIF, "bold": True, "italic": True}])
tabs = []
prices2 = [("Amazon", "₹1,540", False), ("Flipkart", "₹1,610", False), ("Myntra", "₹1,499", True), ("Tata CLiQ", "₹1,585", False)]
x = 0.92
for name, price, win in prices2:
    g = s.shapes.add_group_shape()
    card = rrect(g.shapes, x, 3.5, 2.68, 2.0, fill=CARD, line=(ACCENT if win else LINE), line_w=(2.2 if win else 1.0), radius=0.1)
    add_shadow(card)
    if win:
        add_glow(card, radius_pt=12, alpha=38000)
    for j in range(3):
        oval(g.shapes, x + 0.22 + j * 0.2, 3.72, 0.1, 0.1, fill=FAINT)
    text(g.shapes, x + 0.22, 4.02, 2.3, 0.4, [{"t": name, "size": 15, "color": MUTED, "bold": True}])
    text(g.shapes, x + 0.22, 4.42, 2.3, 0.7, [{"t": price, "size": 33, "color": (ACCENT if win else PAPER), "font": SERIF, "bold": True}])
    tabs.append(g)
    x += 2.93
punch = text(s.shapes, 0.92, 5.95, 11.4, 0.9,
             [{"t": "So which is cheapest? Today you open five tabs and check by hand — and still might overpay.",
               "size": 20, "color": MUTED}], line_spacing=1.2)

tl.fade(eb, 0, 400)
tl.rise(h1, 100); tl.rise(h2, 260)
for i, g in enumerate(tabs):
    tl.rise(g, 550 + i * 150, dur=600, dy=0.045)
tl.rise(punch, 1500)
counts.append(tl.finish(s))

# =====================================================================
# SLIDE 3 — WHAT / IDEA (hub)
# =====================================================================
s = new_slide()
tl = Timeline()
chrome(s, 3)
eb = eyebrow(s, "The idea")
hd = text(s.shapes, 0.9, 1.12, 11.53, 0.9,
          [{"t": "Browse once. We find the best price.", "size": 40, "color": PAPER, "font": SERIF, "bold": True}],
          align=PP_ALIGN.CENTER)
shopper = rrect(s.shapes, 5.39, 2.3, 2.55, 0.68, fill=CARD, line=LINE, radius=0.34)
shape_text(shopper, [{"t": "🧑  A shopper", "size": 17, "color": PAPER, "bold": True}])
ar = arrow(s.shapes, 6.67, 3.06, 6.67, 3.6, color=FAINT, width=2.0)
center = rrect(s.shapes, 4.87, 3.68, 3.6, 1.1, fill=ACCENT, radius=0.18)
set_gradient(center, ACC_HEX, ACC2_HEX, angle=135)
add_glow(center, radius_pt=16, alpha=45000)
set_paras(center.text_frame, [
    {"runs": [{"t": "Outfit Buddy", "size": 30, "color": DARKTX, "font": SERIF, "bold": True}]},
    {"runs": [{"t": "one place to browse", "size": 13, "color": DARKTX, "bold": True}]},
])
lbl = text(s.shapes, 3.5, 4.95, 6.33, 0.45, [{"t": "↓  links you to the cheapest store", "size": 15, "color": FAINT}], align=PP_ALIGN.CENTER)
chips3 = []
total_w = sum(0.56 + len(n) * 0.15 for n, _ in STORES) + 0.24 * 4
x = (13.333 - total_w) / 2
for name, col in STORES:
    w = 0.56 + len(name) * 0.15
    c = rrect(s.shapes, x, 5.6, w, 0.64, fill=CARD, line=LINE, radius=0.5)
    add_shadow(c, blur_pt=8, dist_pt=3, alpha=30000)
    shape_text(c, [{"t": name, "size": 15, "color": col, "bold": True}])
    chips3.append(c)
    x += w + 0.24

tl.fade(eb, 0, 400)
tl.rise(hd, 100)
tl.pop(shopper, 450)
tl.wipe(ar, 720, dur=350, direction="down")
tl.pop(center, 900, dur=600)
tl.fade(lbl, 1280)
for i, c in enumerate(chips3):
    tl.pop(c, 1430 + i * 110)
counts.append(tl.finish(s))

# =====================================================================
# SLIDE 4 — WHAT / VALUE (3 cards)
# =====================================================================
s = new_slide()
tl = Timeline()
chrome(s, 4)
eb = eyebrow(s, "What you get")
hd = text(s.shapes, 0.9, 1.1, 11.5, 0.95, [{"t": "One closet, every store.", "size": 44, "color": PAPER, "font": SERIF, "bold": True}])
vals = [
    ("🛍", "One catalogue", "Browse curated fashion in a single, beautiful place."),
    ("🏷", "The best price", "Every item compared across five stores, automatically."),
    ("⚡", "One tap to buy", "We send you straight to the cheapest checkout."),
]
cards4 = []
x = 0.92
for ic, title, body in vals:
    c = rrect(s.shapes, x, 2.85, 3.72, 3.15, fill=CARD, line=LINE, radius=0.09)
    add_shadow(c)
    set_paras(c.text_frame, [
        {"runs": [{"t": ic, "size": 30, "color": ACCENT}], "align": PP_ALIGN.LEFT},
        {"runs": [{"t": title, "size": 23, "color": PAPER, "font": SERIF, "bold": True}],
         "align": PP_ALIGN.LEFT, "space_before": 16},
        {"runs": [{"t": body, "size": 16, "color": MUTED}],
         "align": PP_ALIGN.LEFT, "space_before": 8, "line_spacing": 1.25},
    ], anchor=MSO_ANCHOR.TOP)
    c.text_frame.margin_left = Inches(0.3)
    c.text_frame.margin_right = Inches(0.3)
    c.text_frame.margin_top = Inches(0.34)
    cards4.append(c)
    x += 3.96

tl.fade(eb, 0, 400)
tl.rise(hd, 100)
for i, c in enumerate(cards4):
    tl.rise(c, 450 + i * 250, dur=650, dy=0.05)
counts.append(tl.finish(s))

# =====================================================================
# SLIDE 5 — HOW / JOURNEY (4 steps)
# =====================================================================
s = new_slide()
tl = Timeline()
chrome(s, 5)
eb = eyebrow(s, "How it works")
hd = text(s.shapes, 0.9, 1.1, 11.6, 0.95, [{"t": "From browse to best price, in four steps.", "size": 40, "color": PAPER, "font": SERIF, "bold": True}])
steps = [("1", "Browse", "Explore a curated fashion catalogue."),
         ("2", "Pick", "Open any product you like."),
         ("3", "Compare", "See its price at every store, side by side."),
         ("4", "Buy", "Tap through to the cheapest one.")]
step_groups, arrows5 = [], []
colw, gap = 2.7, 0.23
x = 0.92
for i, (num, title, body) in enumerate(steps):
    g = s.shapes.add_group_shape()
    cx = x + colw / 2
    ov = oval(g.shapes, cx - 0.45, 3.15, 0.9, 0.9, fill=ACCENT)
    set_gradient(ov, ACC_HEX, ACC2_HEX, angle=135)
    add_glow(ov, radius_pt=10, alpha=35000)
    shape_text(ov, [{"t": num, "size": 27, "color": DARKTX, "font": SERIF, "bold": True}])
    text(g.shapes, x, 4.3, colw, 0.5, [{"t": title, "size": 21, "color": PAPER, "bold": True}], align=PP_ALIGN.CENTER)
    text(g.shapes, x, 4.82, colw, 1.25, [{"t": body, "size": 15, "color": MUTED}], align=PP_ALIGN.CENTER, line_spacing=1.2)
    step_groups.append(g)
    if i < 3:
        arrows5.append(arrow(s.shapes, x + colw - 0.02, 3.6, x + colw + gap + 0.02, 3.6, color=FAINT, width=1.8))
    x += colw + gap

tl.fade(eb, 0, 400)
tl.rise(hd, 100)
for i, g in enumerate(step_groups):
    tl.pop(g, 450 + i * 380, dur=550)
    if i < 3:
        tl.wipe(arrows5[i], 450 + i * 380 + 280, dur=300)
counts.append(tl.finish(s))

# =====================================================================
# SLIDE 6 — HOW / DATA + API CALLS
# =====================================================================
s = new_slide()
tl = Timeline()
chrome(s, 6)
eb = eyebrow(s, "How it works · the data")
hd = text(s.shapes, 0.9, 1.05, 11.6, 0.95, [{"t": "Where the prices come from.", "size": 40, "color": PAPER, "font": SERIF, "bold": True}])
core = rrect(s.shapes, 0.92, 2.8, 2.5, 2.15, fill=CARD2, line=LINE, radius=0.1)
add_shadow(core)
add_glow(core, radius_pt=10, alpha=22000)
set_paras(core.text_frame, [
    {"runs": [{"t": "Outfit Buddy", "size": 30, "color": PAPER, "font": SERIF, "bold": True}]},
    {"runs": [{"t": "runs in your browser", "size": 13, "color": MUTED}], "space_before": 6},
])
api_prices = ["₹1,540", "₹1,610", "₹1,499", "₹1,560", "₹1,585"]
wires, packets, storechips, pricetxts = [], [], [], []
row_y0, row_dy = 2.6, 0.62
for idx, (name, col) in enumerate(STORES):
    y = row_y0 + idx * row_dy
    w = arrow(s.shapes, 3.6, y + 0.16, 8.62, y + 0.16, color=ACCENT, width=2.2)
    wires.append(w)
    pk = oval(s.shapes, 3.72, y + 0.075, 0.17, 0.17, fill=WHITE)
    add_glow(pk, radius_pt=8, alpha=60000)
    packets.append(pk)
    sc = rrect(s.shapes, 8.78, y - 0.07, 1.75, 0.46, fill=CARD, line=LINE, radius=0.3)
    shape_text(sc, [{"t": name, "size": 14, "color": col, "bold": True}])
    storechips.append(sc)
    pt = text(s.shapes, 10.68, y - 0.07, 1.7, 0.46,
              [{"t": api_prices[idx], "size": 19, "color": ACCENT, "font": SERIF, "bold": True}],
              anchor=MSO_ANCHOR.MIDDLE)
    pricetxts.append(pt)
cap = text(s.shapes, 0.92, 5.82, 11.5, 0.55, [
    {"t": "For every product, Outfit Buddy asks ", "size": 19, "color": MUTED},
    {"t": "each store", "size": 19, "color": PAPER, "bold": True},
    {"t": ": “what's your price?” — then collects the answers.", "size": 19, "color": MUTED},
])
foot = text(s.shapes, 0.92, 6.5, 11.5, 0.85,
            [{"t": "In this demo the price answers are simulated (stable & realistic) because live retailer APIs need "
                   "secret keys. The same design plugs into real pricing feeds behind a secure server — with no change to the app.",
              "size": 13, "color": FAINT}], line_spacing=1.2)

tl.fade(eb, 0, 400)
tl.rise(hd, 100)
tl.pop(core, 350, dur=600)
for i in range(5):
    tl.wipe(wires[i], 620 + i * 130, dur=450)
    tl.pop(storechips[i], 800 + i * 130)
    tl.fade(packets[i], 1000 + i * 150, dur=300)
    tl.travel(packets[i], 1150 + i * 150, dx=0.352, dur=1700)
    tl.fade(pricetxts[i], 1650 + i * 140, dur=450)
tl.rise(cap, 2450)
tl.fade(foot, 2850, dur=600)
counts.append(tl.finish(s))

# =====================================================================
# SLIDE 7 — HOW / WHICH IS CHEAPER
# =====================================================================
s = new_slide()
tl = Timeline()
chrome(s, 7)
eb = eyebrow(s, "How it works · the comparison")
hd = text(s.shapes, 0.9, 1.05, 11.6, 0.95, [{"t": "How we know which store is cheapest.", "size": 40, "color": PAPER, "font": SERIF, "bold": True}])
prod = text(s.shapes, 0.92, 2.02, 11, 0.45, [
    {"t": "Comparing one product · ", "size": 17, "color": MUTED},
    {"t": "Linen Overshirt", "size": 19, "color": PAPER, "bold": True},
])
comp = [("Amazon", "₹1,540", "in stock", False, False),
        ("Flipkart", "₹1,610", "in stock", False, False),
        ("Myntra", "₹1,499", "in stock", True, False),
        ("Ajio", "₹1,470", "out of stock", False, True),
        ("Tata CLiQ", "₹1,585", "in stock", False, False)]
cards7 = []
best_card = badge = None
x, cw = 0.92, 2.28
for name, price, stock, best, oos in comp:
    c = rrect(s.shapes, x, 2.72, cw, 2.0, fill=CARD, line=(ACCENT if best else LINE), line_w=(2.4 if best else 1.0), radius=0.1)
    add_shadow(c)
    set_paras(c.text_frame, [
        {"runs": [{"t": name, "size": 15, "color": (FAINT if oos else MUTED), "bold": True}]},
        {"runs": [{"t": price, "size": 30, "color": (ACCENT if best else (FAINT if oos else PAPER)),
                   "font": SERIF, "bold": True, "strike": oos}], "space_before": 6},
        {"runs": [{"t": stock, "size": 13, "color": FAINT}], "space_before": 4},
    ])
    cards7.append(c)
    if best:
        best_card = c
        add_glow(c, radius_pt=14, alpha=42000)
        badge = rrect(s.shapes, x + cw / 2 - 0.85, 2.42, 1.7, 0.44, fill=ACCENT, radius=0.5)
        set_gradient(badge, ACC_HEX, ACC2_HEX, angle=135)
        shape_text(badge, [{"t": "★ Best price", "size": 12, "color": DARKTX, "bold": True}])
    x += cw + 0.19
vt = text(s.shapes, 0.92, 5.25, 6.3, 0.85, [
    {"t": "The lowest price that's actually ", "size": 20, "color": PAPER},
    {"t": "in stock", "size": 20, "color": ACCENT, "bold": True},
    {"t": " wins →", "size": 20, "color": PAPER},
], anchor=MSO_ANCHOR.MIDDLE)
buy = rrect(s.shapes, 7.3, 5.25, 3.6, 0.8, fill=ACCENT, radius=0.5)
set_gradient(buy, ACC_HEX, ACC2_HEX, angle=135)
add_glow(buy, radius_pt=12, alpha=45000)
shape_text(buy, [{"t": "Buy on Myntra · ₹1,499 ↗", "size": 17, "color": DARKTX, "bold": True}])

tl.fade(eb, 0, 400)
tl.rise(hd, 100)
tl.fade(prod, 320)
for i, c in enumerate(cards7):
    tl.pop(c, 480 + i * 130)
tl.pop(badge, 1400, dur=500)
tl.pulse(best_card, 1400)
tl.rise(vt, 1800)
tl.pop(buy, 1980, dur=550)
counts.append(tl.finish(s))

# =====================================================================
# SLIDE 8 — TECHNOLOGY IDEA (adapters)
# =====================================================================
s = new_slide()
tl = Timeline()
chrome(s, 8)
eb = eyebrow(s, "The technology, in one idea")
hd = s.shapes.add_textbox(Inches(0.9), Inches(1.05), Inches(11.7), Inches(1.6))
set_paras(hd.text_frame, [
    {"runs": [{"t": "Every store speaks differently.", "size": 38, "color": PAPER, "font": SERIF, "bold": True}],
     "align": PP_ALIGN.LEFT, "line_spacing": 1.08},
    {"runs": [{"t": "We built a translator for each.", "size": 38, "color": PAPER, "font": SERIF, "bold": True}],
     "align": PP_ALIGN.LEFT, "line_spacing": 1.08},
], anchor=MSO_ANCHOR.TOP)
core8 = rrect(s.shapes, 0.92, 3.3, 2.5, 2.05, fill=ACCENT, radius=0.12)
set_gradient(core8, ACC_HEX, ACC2_HEX, angle=135)
add_glow(core8, radius_pt=14, alpha=42000)
set_paras(core8.text_frame, [
    {"runs": [{"t": "Outfit Buddy", "size": 24, "color": DARKTX, "font": SERIF, "bold": True}]},
    {"runs": [{"t": "core", "size": 24, "color": DARKTX, "font": SERIF, "bold": True}]},
])
plugs = [("Amazon adapter", "?k=linen shirt"), ("Myntra adapter", "/linen-shirt"),
         ("Flipkart adapter", "/search?q=…"), ("…and any new store", "")]
plug_groups, arrows8 = [], []
py = 3.1
for name, code in plugs:
    arrows8.append(arrow(s.shapes, 3.52, py + 0.3, 4.14, py + 0.3, color=FAINT, width=1.8))
    g = s.shapes.add_group_shape()
    pr = rrect(g.shapes, 4.24, py, 6.95, 0.6, fill=CARD, line=LINE, radius=0.2)
    add_shadow(pr, blur_pt=8, dist_pt=3, alpha=30000)
    shape_text(pr, [{"t": "🔌  " + name, "size": 17, "color": PAPER, "bold": True}], align=PP_ALIGN.LEFT)
    pr.text_frame.margin_left = Inches(0.24)
    if code:
        text(g.shapes, 8.3, py + 0.09, 2.68, 0.44,
             [{"t": code, "size": 15, "color": MUTED, "font": SERIF, "italic": True}],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    plug_groups.append(g)
    py += 0.7
cap8 = text(s.shapes, 0.92, 6.15, 11.6, 0.9, [
    {"t": "Each “adapter” hides one store's quirks — so ", "size": 18, "color": MUTED},
    {"t": "adding a new store is a one-line change", "size": 18, "color": PAPER, "bold": True},
    {"t": ", and the rest of the app never notices.", "size": 18, "color": MUTED},
], line_spacing=1.2)

tl.fade(eb, 0, 400)
tl.rise(hd, 100)
tl.pop(core8, 400, dur=600)
for i in range(4):
    tl.wipe(arrows8[i], 620 + i * 170, dur=300)
    tl.rise(plug_groups[i], 700 + i * 170, dur=550, dy=0.035)
tl.rise(cap8, 1650)
counts.append(tl.finish(s))

# =====================================================================
# SLIDE 9 — CLOSE
# =====================================================================
s = new_slide()
tl = Timeline()
r1 = ring(s.shapes, 10.4, -1.4, 3.8)
r2 = ring(s.shapes, -1.5, 5.4, 3.4, alpha=9000)
eb = eyebrow(s, "Outfit Buddy", align=PP_ALIGN.CENTER)
big = s.shapes.add_textbox(Inches(0.9), Inches(1.95), Inches(11.53), Inches(2.7))
set_paras(big.text_frame, [
    {"runs": [{"t": "One closet.", "size": 78, "color": PAPER, "font": SERIF, "bold": True}], "line_spacing": 1.04},
    {"runs": [{"t": "Every store.", "size": 78, "color": PAPER, "font": SERIF, "bold": True}], "line_spacing": 1.04},
], anchor=MSO_ANCHOR.TOP)
sub9 = text(s.shapes, 0.9, 4.95, 11.53, 0.6,
            [{"t": "Browse everywhere. Buy at the best price. Every time.", "size": 22, "color": MUTED}],
            align=PP_ALIGN.CENTER)
chips9 = []
total_w = sum(0.52 + len(n) * 0.145 for n, _ in STORES) + 0.22 * 4
x = (13.333 - total_w) / 2
for name, _c in STORES:
    w = 0.52 + len(name) * 0.145
    c = rrect(s.shapes, x, 5.9, w, 0.55, fill=None, line=LINE, radius=0.5)
    shape_text(c, [{"t": name, "size": 14, "color": MUTED, "bold": True}])
    chips9.append(c)
    x += w + 0.22

tl.fade(r1, 200, 900); tl.fade(r2, 400, 900)
tl.fade(eb, 0, 400)
tl.zoom_slow(big, 150, dur=950)
tl.rise(sub9, 750)
for i, c in enumerate(chips9):
    tl.pop(c, 1000 + i * 100)
counts.append(tl.finish(s))

prs.save(OUT)
print(f"Saved {OUT}")
print("slides:", len(counts), "| animation effects per slide:", counts, "| total:", sum(counts))
